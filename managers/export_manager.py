# 
# *** Export Manager: Manages the data transformation process (eg: from JSON to excel, pptx ...)
from managers.data_manager import DataManager

import collections 
import collections.abc
import subprocess
import sys
import pandas as pd
import requests
import streamlit as st

try:
    from pptx import Presentation
except ModuleNotFoundError:
    print("Module 'pptx' not found. Installing it now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    print("Module 'pptx' installed successfully.")

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import json, os, io, base64

class ExportManager:

    topic_colors = {
    "social": RGBColor(186, 85, 211),
    "technological": RGBColor(30, 144, 255),
    "environmental": RGBColor(10, 160, 160),
    "economic": RGBColor(220, 20, 60),
    "political": RGBColor(139, 69, 19),
    "business_and_investment": RGBColor(255, 165, 0)
    }

    topic_titles = {
        "social": "social",
        "technological": "technological",
        "environmental": "environmental",
        "economic": "economic",
        "political": "political",
        "business_and_investment": "business_and_investment"
    }

    # *** PPTX Utilities ***
    @staticmethod
    def init_slides():
        prs = Presentation()
        prs.slide_width = Inches(16.0)  
        prs.slide_height = Inches(9.0)  
        return prs

    # Define function to add slides with optional fixed font size
    @staticmethod
    def add_slide(prs, topic, title, content, title_font_size, content_font_size, background_color=RGBColor(240, 240, 240), fixed_font_size=False, rgb_color = "#F3B915"):
        slide_layout = prs.slide_layouts[1]  # Use title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]

        # Adjust title placeholder size to avoid overlapping with the black line
        title_placeholder.width = Inches(12.5)
        title_placeholder.height = Inches(1.2)
        title_placeholder.left = Inches(1)
        title_placeholder.top = Inches(0.4)

        # Set title and font size
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.clear()
        title_p = title_text_frame.add_paragraph()
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.bold = True
        title_p.font.name = 'Microsoft JhengHei'
        title_p.font.color.rgb = ExportManager.topic_colors.get(topic, rgb_color)
        title_p.alignment = PP_ALIGN.LEFT

        # Add a black line below the title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.9), Inches(14), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.width = Inches(0)

        # Add solid circle in the upper-right corner
        solid_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(13.7), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        solid_circle.fill.solid()
        solid_circle.fill.fore_color.rgb = ExportManager.topic_colors.get(topic, rgb_color)

        # Add hollow circle in the upper-right corner
        hollow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(14.5), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        hollow_circle.line.color.rgb = RGBColor(0, 0, 0)
        hollow_circle.line.width = Pt(1)
        hollow_circle.fill.background()

        # Adjust content placeholder size
        body_placeholder.width = Inches(14)
        body_placeholder.height = Inches(6)
        body_placeholder.left = Inches(1)
        body_placeholder.top = Inches(2.05)

        # Set content and font size
        text_frame = body_placeholder.text_frame
        text_frame.clear()

        # Check if font size should be fixed or auto-adjusted
        if fixed_font_size:
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-sizing
        else:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Enable auto-sizing

        paragraphs = content.split('\n')
        for paragraph in paragraphs:
            p = text_frame.add_paragraph()
            parts = paragraph.split('**')
            for i, part in enumerate(parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(content_font_size)
                run.font.name = 'Microsoft JhengHei'
                if i % 2 == 1:  # Odd index means text between asterisks
                    run.font.bold = True

    # Use this function with `fixed_font_size=True` for Trend Insights and Hashtag slides
    @staticmethod
    def add_report_and_keywords(prs, data, topic, rgb_color):
        # For the Trend Insights and Hashtag slide, use fixed font size
        # todo *************
        slide_title = ExportManager.topic_titles.get(topic, topic)
        ExportManager.add_slide(prs = prs,
                title = slide_title, 
                content = f"{data['趨勢報告']}\n\n**關鍵字：**\n" + "、".join(data["關鍵字"]),
                title_font_size=32, content_font_size=26, fixed_font_size=False, topic = topic, rgb_color = rgb_color)
        
        # For the other slide, allow auto-adjustment
        second_page = ""
        for i in range(1, 13):
            trend_key = f"主要趨勢{i}"
            if trend_key in data:  # 檢查主要趨勢是否存在
                second_page += data[trend_key].get("標題", "") + "\n\n"
            else:
                # print(f"{trend_key} 不存在，跳過該趨勢的處理。")
                pass


        ExportManager.add_slide(prs = prs, topic = topic, title = slide_title, content = second_page, title_font_size=32, content_font_size=26, rgb_color = rgb_color)  # Allow auto-adjustment for this slide

    # Define function to add main trend pages (for STEEP)
    @staticmethod
    def add_trend_pages(prs, data, trend_number, trend_type, topic, rgb_color):
        trend_key = f"{trend_type}{trend_number}"
        if trend_key in data:
            trend = data[trend_key]

            # First Slide: Trend Insights and Hashtag Keywords
            ExportManager.add_slide(prs = prs,
                topic = topic,
                title = f"{trend_type}{trend_number}： {trend['標題']}",
                content = f"**趨勢洞察：**{trend['<a>趨勢洞察']}\n\n**Hashtag關鍵詞：**\n" + 
                "、".join(str(hashtag) for hashtag in trend["<b>Hashtag關鍵詞"]),
                title_font_size=32, content_font_size=22, fixed_font_size=True # Adjust content font size to 26
                , rgb_color = rgb_color
            )

            # Second Slide: Representative Events
            events_content = "\n\n".join(
                f"**事件：**{event.get('事件', '無資料')}\n**分析：**{event.get('分析', '無資料')}\n**來源：**{event.get('(來源', '無資料')}\n**關聯度：**{event.get('關聯度', '無資料')}"
                if isinstance(event, dict) else str(event)
                for event in trend["<c>代表事件"]
            )
            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}： {trend['標題']}",
                content = events_content,
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

            # Third Slide: Key Stakeholders, Issue Gaps, and Future Service Opportunities
            important_stakeholders = "\n".join([f"{key}: {value}" for key, value in trend["<d>重要關係人"].items()])
            issue_gaps = "\n".join(str(gap) for gap in trend["<e>缺口"])
            
            # 檢查 "<f>未來產品或服務機會點" 是否存在
            future_opportunities = "\n".join(str(opportunity) for opportunity in trend.get("<f>未來產品或服務機會點", []))
            
            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}：{trend['標題']} ",
                content = f"**重要關係人：**\n{important_stakeholders}\n\n**議題缺口：**\n{issue_gaps}\n\n**未來服務機會點：**\n{future_opportunities}",
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

            # Fourth Slide: Key Drivers, Weak Signals, and Trend Summary Insights
            # key_drivers = "\n".join([f"{key}: {value}" for key, value in trend["<g>關鍵驅動因素"].items()])
            key_drivers = "\n".join(["：".join([str(aspect), str(driver)])
                                    for aspect, driver in trend["<g>關鍵驅動因素"].items()])
            weak_signals = "\n".join([str(signal) for signal in trend["<h>微弱信號"]])
            time_scale = trend.get("<i>時間尺度", "")
            trend_summary = trend.get("<j>趨勢總結洞察", "")
            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}：{trend['標題']} ",
                content = f"**關鍵驅動因素：**\n{key_drivers}\n\n**微弱信號：**\n{weak_signals}\n\n**時間尺度：**{time_scale}\n\n**趨勢總結洞察：**\n{trend_summary}",
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

    # Define functions to add trend pages flexibly(for SELF SELECT)
    @staticmethod
    def add_trend_pages_flexibly(prs, data, trend_number, trend_type, topic, rgb_color):
        trend_key = f"{trend_type}{trend_number}"
        if trend_key in data:
            trend = data[trend_key]

            # First Slide: Trend Insights and Hashtag Keywords
            ExportManager.add_slide(prs = prs,
                topic = topic,
                title = f"{trend_type}{trend_number}： {trend['標題']}",
                content = f"**趨勢洞察：**{trend['<a>趨勢洞察']}\n\n**Hashtag關鍵詞：**\n" + 
                "、".join(str(hashtag) for hashtag in trend["<b>Hashtag關鍵詞"]),
                title_font_size=32, content_font_size=22, fixed_font_size=True # Adjust content font size to 26
                , rgb_color = rgb_color
            )

            # Second Slide: Representative Events
            events_content = "\n\n".join(
                f"**事件：**{event.get('事件', '無資料')}\n**分析：**{event.get('分析', '無資料')}\n**來源：**{event.get('(來源', '無資料')}\n**關聯度：**{event.get('關聯度', '無資料')}"
                if isinstance(event, dict) else str(event)
                for event in trend["<c>代表事件"]
            )
            ExportManager.add_slide(prs = prs, topic = topic,
                title = f"{trend_type}{trend_number}： {trend['標題']}",
                content = events_content,
                title_font_size=32, content_font_size=28, rgb_color = rgb_color
            )

            # *** Other Inference Results ***
            corpus = []

            # 重要關係人
            try:
                important_stakeholders = "\n".join([f"{key}: {value}" for key, value in trend["<d>重要關係人"].items()])
                corpus.append({"**重要關係人**": important_stakeholders})
            except:
                pass
            # 議題缺口
            try:
                issue_gaps = "\n".join(str(gap) for gap in trend["<e>缺口"])
                corpus.append({"**議題缺口**": issue_gaps})
            except:
                pass
            # 議題缺口
            try:
                future_opportunities = "\n".join(str(opportunity) for opportunity in trend["<f>未來產品或服務機會點"])
                corpus.append({"**未來服務機會點**": future_opportunities})
            except:
                pass
            # 關鍵驅動因素
            try:
                key_drivers = "\n".join(["：".join([str(aspect), str(driver)]) for aspect, driver in trend["<g>關鍵驅動因素"].items()])
                corpus.append({"**關鍵驅動因素**": key_drivers})
            except:
                pass
            # 微弱信號
            try:
                weak_signals = "\n".join([str(signal) for signal in trend["<h>微弱信號"]])
                corpus.append({"**微弱信號**": weak_signals})
            except:
                pass
            # 時間尺度
            try:
                time_scale = trend["<i>時間尺度"]
                corpus.append({"**時間尺度**": time_scale})
            except:
                pass
            # 趨勢總結洞察
            try:
                trend_summary = trend["<j>趨勢總結洞察"]
                corpus.append({"**趨勢總結洞察**": trend_summary})
            except:
                pass
            
            if len(corpus) > 3:
                chunk_size = 3
                chunks = [corpus[i:i + chunk_size] for i in range(0, len(corpus), chunk_size)]
            else:
                chunks = [corpus]
            
            for chunk in chunks:
                chunk_dict = {}
                for _ in chunk:
                    chunk_dict.update(_)
                ExportManager.add_slide(prs = prs, topic = topic,
                    title = f"{trend_type}{trend_number}：{trend['標題']} ",
                    content = "\n\n".join([f"{key}:\n{value}\n\n" for key, value in chunk_dict.items()]),
                    title_font_size=32, content_font_size=28, rgb_color = rgb_color
                )

    
    # Function that queries pics with project name as keyword, and downloads queried pictures
    @staticmethod
    def download_images(query, num_images):
        # * query
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": st.secrets["PEXEL_KEY"]}
        params = {"query": query, "per_page": num_images}
        response = requests.get(url, headers=headers, params=params)
        response.raise_for_status()
        results = response.json()
        image_urls = [photo["src"]["original"] for photo in results["photos"]]

        os.makedirs("images", exist_ok = True)
        for i, url in enumerate(image_urls):
            response = requests.get(url)
            with open(f"images/image_{i+1}.jpg", "wb") as f:
                f.write(response.content)
    
    # Function that add downloaded pics to ppt
    @staticmethod
    def add_image_slide(prs, image_file, title, title_font_size, background_color = RGBColor(240, 240, 240) , fixed_font_size=False, rgb_color = RGBColor(5, 5, 5),
                        ):
        slide_layout = prs.slide_layouts[0]  # Use title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title_placeholder = slide.shapes.title
        

        # Adjust title placeholder size to avoid overlapping with the black line
        title_placeholder.width = Inches(12.5)
        title_placeholder.height = Inches(1.2)
        title_placeholder.left = Inches(1)
        title_placeholder.top = Inches(0.4)

        # Set title and font size
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.clear()
        title_p = title_text_frame.add_paragraph()
        title_p.text = title
        title_p.font.size = Pt(title_font_size)
        title_p.font.bold = True
        title_p.font.name = 'Microsoft JhengHei'
        title_p.font.color.rgb = rgb_color
        title_p.alignment = PP_ALIGN.LEFT

        # Add a black line below the title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.9), Inches(14), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0, 0, 0)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.width = Inches(0)

        # Add solid circle in the upper-right corner
        solid_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(13.7), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        solid_circle.fill.solid()
        solid_circle.fill.fore_color.rgb = rgb_color

        # Add hollow circle in the upper-right corner
        hollow_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(14.5), Inches(0.9), Inches(0.8), Inches(0.8)
        )
        hollow_circle.line.color.rgb = RGBColor(0, 0, 0)
        hollow_circle.line.width = Pt(1)
        hollow_circle.fill.background()

    

        # Add pics

        image_path = os.path.join("images", image_file)
        slide.shapes.add_picture(image_path, Inches(4), Inches(2), Inches(8), Inches(6))  # 調整圖片大小和位置


    # todo *** EXCEL Utilities ***
    # Transform the final output json data into pd.DataFrame
    @staticmethod
    def get_report_excels(start_date: str, end_date: str, topic: str, data = None) -> pd.DataFrame:
        
        if data == None:
            data = DataManager.get_files(f"{topic}_trend_report_{start_date}-{end_date}.json", 'json')
            data = json.loads(base64.b64decode(data))
        # data: {'趨勢報告': '...', '關鍵字': [], '主要趨勢i': {'趨勢洞察': '...', '關鍵字': '...',....}, ...}
        # '趨勢報告'，'關鍵字' / 主要趨勢i 存成兩個 dataframe（要分開擺放在 Excel 中的不同 worksheet）

        # ***************
        # *** 主要趨勢 ***
        # ** excel_df1 **

        # get keys from the json data
        json_keys = {}
        
        for key, value in data.items():

            if isinstance(data[key], dict):
                cols = []
                for key_1, value_1 in data[key].items():
                    cols.append(key_1)

                json_keys[key] = cols

            else:
                pass

        # initialize dataframe
        excel_df1 = pd.DataFrame(columns = json_keys['主要趨勢1'])

        # 將「是字典的值取出」（主要趨勢）並放入 dataframe
        for trend, aspects in data.items():

            if isinstance(aspects, dict):

                row = []

                for title, aspect in aspects.items():
                    row.append(aspect)

                row = pd.DataFrame(row).T
                row.columns = json_keys['主要趨勢1']
        
                excel_df1 = pd.concat([excel_df1, row], axis = 0)

        # 處理各欄位的文字資料
                
            # Hashtag
        try:   
            for i, item in enumerate(excel_df1.loc[:, '<b>Hashtag關鍵詞']):
                excel_df1.iloc[i, 2] = '\n'.join(item)
        except:
            pass

            # 事件
        try:
            for i, item in enumerate(excel_df1.loc[:, '<c>代表事件']):
                text = ""
                for j, value in enumerate(item):
                    try:
                        source = value['來源']
                    except:
                        try:
                            source = value['(來源']
                        except:
                            source = None

                    str = "\n".join([f"事件{j+1}: {value['事件']}", f"分析: {value['分析']}", f"來源: {source}", f"關聯度: {value['關聯度']}", "******\n"])
                    text += str
                excel_df1.iloc[i, 3] = text
        except:
            pass

            # 關係人
        try:
            for i, item in enumerate(excel_df1.loc[:, '<d>重要關係人']):
                text = ""
                text += "\n".join([f"供給端: {item['供給端']}", f"需求端: {item['需求端']}", f"代表性意義: {item['代表性意義']}"])
                excel_df1.iloc[i, 4] = text
        except:
            pass

            # 缺口
        try:
            for i, item in enumerate(excel_df1.loc[:, '<e>缺口']):
                text = "\n".join(item)
                excel_df1.iloc[i, 5] = text
        except:
            pass

            # 產品服務
        try:
            for i, item in enumerate(excel_df1.loc[:, '<f>未來產品或服務機會點']):
                text = "\n".join(item)
                excel_df1.iloc[i, 6] = text
        except:
            pass

            # 驅動因素
        try:
            for i, item in enumerate(excel_df1.loc[:, '<g>關鍵驅動因素']):
                text = ""

                try:
                    A = item['a社會']
                except:
                    A = None

                try:
                    B = item['b政治']
                except:
                    B = None

                try:
                    C = item['c經濟']
                except:
                    C = None

                try:
                    D = item['d文化']
                except:
                    D = None


                text += "\n".join([f"社會: {A}", f"政治: {B}", f"經濟: {C}", f"文化: {D}"])
                excel_df1.iloc[i, 7] = text
        except:
            pass

            # 微弱信號
        try:
            for i, item in enumerate(excel_df1.loc[:, '<h>微弱信號']):
                text = "\n".join(item)
                excel_df1.iloc[i, 8] = text
        except:
            pass


        # *********************
        # *** 趨勢洞察、關鍵字 ***
        # ***** excel_df2 *****

        excel_df2 = pd.DataFrame(data = {'趨勢報告': data['趨勢報告'], '關鍵字': "\n".join(data['關鍵字'])}, index = [topic])

        return excel_df1, excel_df2
                
    class STEEP:
        
        @staticmethod
        def create_pptx(topic, data = None):
            
            color = '#808080'
            rgb_color = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))

            prs = ExportManager.init_slides()
            assert data is not None, "Input data is required to create slides!"
            ExportManager.add_report_and_keywords(prs, data, topic, rgb_color)
            
            for i in range(1, 13):
                try:
                    ExportManager.add_trend_pages(prs, data, i, "主要趨勢", topic, rgb_color)
                except:
                    pass

            


            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            base64_pptx = base64.b64encode(buffer.read()).decode('utf-8')
            buffer.close()

            return base64_pptx

        @staticmethod
        def create_excel(start_date, end_date, topics: list):
            
            summary_reports = []

            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer, engine = 'openpyxl') as writer:
                for topic in topics:
                    try:
                        dfs_to_write = ExportManager.get_report_excels(start_date, end_date, topic)
                        summary_reports.append(dfs_to_write[1])
                        dfs_to_write[0].to_excel(writer, sheet_name = topic, index = False)
                    except:
                        pass
                agg_res = pd.concat(summary_reports, ignore_index = False)
                agg_res.to_excel(writer, sheet_name = 'STEEP 趨勢報告與關鍵字', index = True)

            buffer.seek(0)
            b64_excel = base64.b64encode(buffer.read()).decode()
            buffer.close()
            

            return b64_excel
        
    class SELF_SELECT:
        
        @staticmethod
        def create_pptx(color, topic, data = None):

            # * transform the html color code to RGBColor object
            rgb_color = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))

            prs = ExportManager.init_slides()
            assert data is not None, "Input data is required to create slides!"

            ExportManager.add_report_and_keywords(prs, data, topic, rgb_color)

            for i in range(1, 13):
                try:
                    ExportManager.add_trend_pages_flexibly(prs, data, i, "主要趨勢", topic, rgb_color)
                except:
                    pass
                i += 1
            
            ExportManager.download_images(topic, 5)
            for i in range(1, 6):
                ExportManager.add_image_slide(
                    prs, f"image_{i}.jpg", 'TEST', 30
                )

            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            base64_pptx = base64.b64encode(buffer.read()).decode('utf-8')
            buffer.close()

            return base64_pptx
        
        @staticmethod
        def create_excel(start_date, end_date, title):
            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer) as writer:
                output = ExportManager.get_report_excels(start_date, end_date, title)
                output[1].to_excel(writer, sheet_name = '趨勢報告與關鍵字', index = False)
                output[0].to_excel(writer, sheet_name = '細項', index = False)
            buffer.seek(0)
            b64_excel = base64.b64encode(buffer.read()).decode()
            buffer.close()

            return b64_excel





# # *** Merge all dataframes of STEEP trend report into one excel file ***
# def steep_merge_as_one_Excel(start_time: str, end_time: str, topics: list, user_name: str, user_email: str):
#     summary_report = []
#     path = f'./output/{start_time}-{end_time}_STEEP.xlsx'
#     with pd.ExcelWriter(path) as writer:
#         for topic in topics:
#             try:
#                 summary_report.append(json_to_df(start_time, end_time, topic)[1])
#                 json_to_df(start_time, end_time, topic)[0].to_excel(writer, sheet_name = topic, index = False)
#             except:
#                 pass
#         summary_report = pd.concat(summary_report, ignore_index = False)
#         summary_report.to_excel(writer, sheet_name = 'STEEP 趨勢報告與關鍵字', index = True)

#     with open(rf'./output/{start_time}-{end_time}_STEEP.xlsx', "rb") as file:
#         contents = file.read()
#         b64_excel = base64.b64encode(contents).decode()
#     response = postFiles(file_name = f"{start_time}-{end_time}_STEEP.xlsx", file_content = b64_excel, 
#                     expiration = str(dt.datetime.today() + dt.timedelta(90)), user_name = user_name, user_email = user_email)


#     return response

# # *** SELF_SELECT_PAGE() Export as Excel ***
# def self_select_to_Excel(start_date: str, end_date: str, title: str, user_name: str, user_email: str):
#     path = f'./output/{title}_{start_date}-{end_date}_trend_report.xlsx'
#     with pd.ExcelWriter(path) as writer:
#         output = json_to_df(start_date, end_date, title)
#         output[1].to_excel(writer, sheet_name = '趨勢報告與關鍵字', index = False)
#         output[0].to_excel(writer, sheet_name = '細項', index = False)

#     with open(rf'./output/{title}_{start_date}-{end_date}_trend_report.xlsx', "rb") as file:
#         contents = file.read()
#         b64_excel = base64.b64encode(contents).decode()
#     response = postFiles(file_name = f"{title}_{start_date}-{end_date}_trend_report.xlsx", file_content = b64_excel, 
#                     expiration = str(dt.datetime.today() + dt.timedelta(90)), user_name = user_name, user_email = user_email)

